#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT翻訳のコア機能
"""

import os
import re
import time
from pathlib import Path
from typing import Callable, Optional

from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from ..utils.config import ConfigManager

# OCI SDK for Grok-3 support
try:
    import oci
    OCI_AVAILABLE = True
except ImportError:
    OCI_AVAILABLE = False
    print("警告: OCI SDK が見つかりません。Grok-3 モデルを使用するには 'pip install oci' でインストールしてください。")


class PPTTranslator:
    """PPT翻訳クラス"""

    def __init__(self):
        self.config_manager = ConfigManager()
        self.openai_client = None
        self.oci_client = None
        self._init_clients()

    def _init_clients(self):
        """各種クライアントの初期化"""
        # OpenAI クライアント
        api_key = self.config_manager.get_openai_api_key()
        base_url = self.config_manager.get_openai_base_url()

        if api_key and base_url:
            self.openai_client = OpenAI(
                api_key=api_key,
                base_url=base_url
            )

        # OCI クライアント (Grok-3用)
        if OCI_AVAILABLE:
            try:
                # OCI設定を読み込み
                config_profile = self.config_manager.get_config_profile()
                config = oci.config.from_file('~/.oci/config', config_profile)
                endpoint = "https://inference.generativeai.us-chicago-1.oci.oraclecloud.com"
                self.oci_client = oci.generative_ai_inference.GenerativeAiInferenceClient(
                    config=config,
                    service_endpoint=endpoint,
                    retry_strategy=oci.retry.NoneRetryStrategy(),
                    timeout=(10, 240)
                )
                # 設定からcompartment_idを取得
                self.compartment_id = self.config_manager.get_compartment_id()
                if not self.compartment_id:
                    raise Exception("COMPARTMENT_ID が設定されていません。設定ダイアログでOCI Compartment IDを設定してください。")
            except Exception as e:
                print(f"OCI クライアントの初期化に失敗しました: {e}")
                self.oci_client = None
    
    def translate_text(self, text: str, target_lang: str, model_name: str = "gpt-4o") -> str:
        """テキストを翻訳"""
        max_attempts = 5  # 最大試行回数

        # モデルに応じて適切な翻訳方法を選択
        if model_name.startswith("xai.grok"):
            return self._translate_with_grok(text, target_lang, model_name, max_attempts)
        else:
            return self._translate_with_openai(text, target_lang, model_name, max_attempts)

    def _translate_with_openai(self, text: str, target_lang: str, model_name: str, max_attempts: int) -> str:
        """OpenAI APIを使用して翻訳"""
        if not self.openai_client:
            raise Exception("OpenAIクライアントが初期化されていません。API設定を確認してください。")

        for attempt in range(max_attempts):
            try:
                completion = self.openai_client.chat.completions.create(
                    model=model_name,
                    messages=[
                        {
                            "role": "system",
                            "content": "あなたはプレゼンテーションスライドで使用される簡潔で専門的なスタイルを維持することに特化した翻訳アシスタントです。プレースホルダーを保持し、言語固有のスタイル規則を尊重しながらテキストを翻訳することがあなたの仕事です。"
                        },
                        {
                            "role": "user",
                            "content": f"""
以下のテキストを{target_lang}に翻訳してください。以下のルールに従ってください：

1. 元のトーンとスタイルを保ち、翻訳がプレゼンテーションスライドに適した簡潔なものになるようにしてください。
2. 過度に正式または冗長な表現は避けてください。例えば：
   - 日本語では、絶対に必要でない限り「です」「ます」を避けてください。
   - 中国語では、率直で専門的な表現を使用してください。
   - 英語では、簡潔さと明確さを優先してください。
3. [PLACEHOLDER_X]形式のプレースホルダーを翻訳または変更しないでください。プレースホルダーは元の位置に、元のテキストと同じ数量で残してください。
4. 説明や追加のコメントなしに、翻訳されたテキストのみを出力してください。

テキスト: {text}
"""
                        }
                    ]
                )

                translated_text = completion.choices[0].message.content
                print(f"原文: {text}")
                print(f"翻訳: {translated_text}\n")
                return translated_text

            except Exception as e:
                if attempt < max_attempts - 1:
                    print(f"試行 {attempt + 1} がエラーで失敗しました: {str(e)}. 再試行中...")
                    time.sleep(1)
                else:
                    print(f"全ての {max_attempts} 回の試行がエラーで失敗しました: {str(e)}")
                    return text

    def _translate_with_grok(self, text: str, target_lang: str, model_name: str, max_attempts: int) -> str:
        """OCI Grok-3 APIを使用して翻訳"""
        if not self.oci_client:
            raise Exception("OCI クライアントが初期化されていません。OCI設定を確認してください。")

        for attempt in range(max_attempts):
            try:
                # プロンプトを作成
                prompt = f"""あなたはプレゼンテーションスライドで使用される簡潔で専門的なスタイルを維持することに特化した翻訳アシスタントです。

以下のテキストを{target_lang}に翻訳してください。以下のルールに従ってください：

1. 元のトーンとスタイルを保ち、翻訳がプレゼンテーションスライドに適した簡潔なものになるようにしてください。
2. 過度に正式または冗長な表現は避けてください。例えば：
   - 日本語では、絶対に必要でない限り「です」「ます」を避けてください。
   - 中国語では、率直で専門的な表現を使用してください。
   - 英語では、簡潔さと明確さを優先してください。
3. [PLACEHOLDER_X]形式のプレースホルダーを翻訳または変更しないでください。プレースホルダーは元の位置に、元のテキストと同じ数量で残してください。
4. 説明や追加のコメントなしに、翻訳されたテキストのみを出力してください。

テキスト: {text}"""

                # OCI Grok-3 リクエストを作成
                chat_detail = oci.generative_ai_inference.models.ChatDetails()

                content = oci.generative_ai_inference.models.TextContent()
                content.text = prompt
                message = oci.generative_ai_inference.models.Message()
                message.role = "USER"
                message.content = [content]

                chat_request = oci.generative_ai_inference.models.GenericChatRequest()
                chat_request.api_format = oci.generative_ai_inference.models.BaseChatRequest.API_FORMAT_GENERIC
                chat_request.messages = [message]
                chat_request.max_tokens = 600
                chat_request.temperature = 1
                chat_request.frequency_penalty = 0
                chat_request.presence_penalty = 0
                chat_request.top_p = 1
                chat_request.top_k = 0

                chat_detail.serving_mode = oci.generative_ai_inference.models.OnDemandServingMode(model_id=model_name)
                chat_detail.chat_request = chat_request
                chat_detail.compartment_id = self.compartment_id

                # API呼び出し
                chat_response = self.oci_client.chat(chat_detail)

                # レスポンスから翻訳結果を抽出
                translated_text = chat_response.data.chat_response.choices[0].message.content[0].text
                print(f"原文: {text}")
                print(f"翻訳: {translated_text}\n")
                return translated_text

            except Exception as e:
                if attempt < max_attempts - 1:
                    print(f"試行 {attempt + 1} がエラーで失敗しました: {str(e)}. 再試行中...")
                    time.sleep(1)
                else:
                    print(f"全ての {max_attempts} 回の試行がエラーで失敗しました: {str(e)}")
                    return text

    def translate_ppt(
        self,
        model_name: str,
        input_ppt: str,
        target_lang: str,
        progress_callback: Optional[Callable[[int], None]] = None,
        status_callback: Optional[Callable[[str], None]] = None,
        log_callback: Optional[Callable[[str], None]] = None,
        stop_callback: Optional[Callable[[], bool]] = None
    ) -> str:
        """PPTファイルを翻訳"""

        def log(message: str):
            """ログ出力"""
            if log_callback:
                log_callback(message)
            print(message)

        def update_status(message: str):
            """ステータス更新"""
            if status_callback:
                status_callback(message)
            log(message)

        def update_progress(value: int):
            """進捗更新"""
            if progress_callback:
                progress_callback(min(100, max(0, value)))  # 0-100の範囲に制限
        
        try:
            # 入力ファイルの確認
            if not os.path.exists(input_ppt):
                raise FileNotFoundError(f"入力ファイルが見つかりません: {input_ppt}")
            
            # 出力ディレクトリの確保
            output_dir = self.config_manager.get_output_directory()
            if not output_dir:
                output_dir = os.path.join(os.path.dirname(input_ppt), "outputs")
            
            os.makedirs(output_dir, exist_ok=True)
            
            # PPTファイルを読み込み
            update_status("PPTファイルを読み込み中...")
            update_progress(5)  # ファイル読み込み開始

            ppt = Presentation(input_ppt)
            input_file_name = os.path.basename(input_ppt)

            total_slides = len(ppt.slides)
            log(f"総スライド数: {total_slides}")
            update_progress(10)  # ファイル読み込み完了

            # 翻訳対象のテキスト要素を事前にカウント（より正確な進捗計算のため）
            update_status("翻訳対象を分析中...")
            total_text_elements = 0
            slide_text_counts = []

            for slide in ppt.slides:
                slide_text_count = 0
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text and run.text.strip():
                                    slide_text_count += 1
                    elif shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text_frame.text and cell.text_frame.text.strip():
                                    slide_text_count += 1

                # ノートも含める
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    slide_text_count += 1

                slide_text_counts.append(slide_text_count)
                total_text_elements += slide_text_count

            log(f"総翻訳対象テキスト要素数: {total_text_elements}")
            update_progress(15)  # 分析完了

            # 各スライドを翻訳
            processed_text_elements = 0

            for slide_index, slide in enumerate(ppt.slides, start=1):
                # 停止チェック
                if stop_callback and stop_callback():
                    log("翻訳が停止されました")
                    update_status("翻訳が停止されました")
                    raise Exception("翻訳が停止されました")

                update_status(f"スライド {slide_index}/{total_slides} を翻訳中...")
                log(f'スライド {slide_index}/{total_slides} を翻訳中')
                log('-------------------------------------------')

                # 基本進捗（15%から85%の範囲で計算）
                base_progress = 15 + int((slide_index - 1) / total_slides * 70)
                update_progress(base_progress)
                
                # スライド内の図形を処理
                for shape in slide.shapes:
                    # フッター部分をスキップ
                    if shape.is_placeholder and shape.placeholder_format.type in [
                        PP_PLACEHOLDER_TYPE.FOOTER,
                        PP_PLACEHOLDER_TYPE.SLIDE_NUMBER,
                        PP_PLACEHOLDER_TYPE.DATE,
                    ]:
                        continue

                    # テーブルの処理
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                original_text = cell.text_frame.text
                                if original_text and original_text.strip() and len(original_text.strip()) > 0:
                                    # 数字かどうかを判定（整数、負数、小数）
                                    if re.match(r'^-?\d+\.?\d*$', original_text.strip()):
                                        continue
                                    translated_text = self.translate_text(original_text, target_lang, model_name)
                                    cell.text_frame.text = translated_text

                                    # 進捗更新（より細かく）
                                    processed_text_elements += 1
                                    if total_text_elements > 0:
                                        detailed_progress = 15 + int((processed_text_elements / total_text_elements) * 70)
                                        update_progress(detailed_progress)
                    
                    # テキストフレームの処理
                    elif shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            # ステップ1: 段落の完全なテキストを抽出し、各runに一意の区切り文字を追加
                            original_runs = []
                            full_text_with_delimiters = ""

                            for idx, run in enumerate(paragraph.runs):
                                original_text = run.text.strip()
                                if original_text and len(original_text) > 0:
                                    delimiter = f"[PLACEHOLDER_{idx}]"  # 一意の区切り文字
                                    full_text_with_delimiters += f"{delimiter}{original_text}"
                                    original_runs.append({"run": run, "delimiter": delimiter})

                            if full_text_with_delimiters == "" or full_text_with_delimiters.strip() == "":
                                continue

                            # 数字かどうかを判定（整数、負数、小数）
                            if re.match(r'^-?\d+\.?\d*$', full_text_with_delimiters.strip()):
                                continue

                            # ステップ2: 段落全体を翻訳（区切り文字を含む）
                            translated_text_with_delimiters = self.translate_text(
                                full_text_with_delimiters, target_lang, model_name
                            )

                            # ステップ3: 区切り文字に基づいて翻訳結果を分割し、各runに書き戻し
                            for item in original_runs:
                                delimiter = item["delimiter"]
                                run = item["run"]

                                # 区切り文字の位置を見つけ、対応する翻訳テキストを抽出
                                start_idx = translated_text_with_delimiters.find(delimiter)
                                if start_idx != -1:
                                    end_idx = start_idx + len(delimiter)
                                    # 翻訳後の内容を抽出し、区切り文字を除去
                                    translated_run_text = translated_text_with_delimiters[end_idx:].split("[PLACEHOLDER_", 1)[0]
                                    run.text = translated_run_text

                            # 進捗更新（段落単位で）
                            processed_text_elements += len(original_runs)
                            if total_text_elements > 0:
                                detailed_progress = 15 + int((processed_text_elements / total_text_elements) * 70)
                                update_progress(detailed_progress)
                
                # ノートスライドの処理
                if slide.has_notes_slide:
                    original_text = slide.notes_slide.notes_text_frame.text
                    if original_text and original_text.strip() and len(original_text.strip()) > 0:
                        # 数字かどうかを判定（整数、負数、小数）
                        if re.match(r'^-?\d+\.?\d*$', original_text.strip()):
                            continue
                        translated_text = self.translate_text(original_text, target_lang, model_name)
                        slide.notes_slide.notes_text_frame.text = translated_text

                        # 進捗更新（ノート翻訳後）
                        processed_text_elements += 1
                        if total_text_elements > 0:
                            detailed_progress = 15 + int((processed_text_elements / total_text_elements) * 70)
                            update_progress(detailed_progress)
            
            # 翻訳後のPPTを保存
            update_status("翻訳結果を保存中...")
            update_progress(90)

            # 元のファイルの拡張子を保持
            original_extension = Path(input_file_name).suffix
            output_file_name = f"{Path(input_file_name).stem}_{target_lang}{original_extension}"
            output_file_path = os.path.join(output_dir, output_file_name)

            log(f"出力ファイルパス: {output_file_path}")
            update_progress(95)

            ppt.save(output_file_path)
            update_progress(98)

            update_status("翻訳が完了しました")
            log("翻訳が完了しました。")
            update_progress(100)
            
            return output_file_path
            
        except Exception as e:
            error_msg = f"翻訳中にエラーが発生しました: {str(e)}"
            log(error_msg)
            raise Exception(error_msg)
